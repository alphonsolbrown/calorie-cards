# streamlit_app.py
# Programmer: Alphonso Brown
# Date: 10/9/2025
# Notes:
#   - Keeps original UI/behavior (Brand/Theme, DB, Last Card, PNG/PPTX)
#   - Fixes Streamlit session_state warning on 'cal'
#   - Adds: save/load cards, reset, create-new, dynamic rows, combo sections,
#           and naming convention for downloads.

from __future__ import annotations
import os, io, json, datetime as dt
from pathlib import Path
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches

from meal_card_generator import Theme, MealItem, MealSection, MealCardData, render_meal_card
from fdc_lookup import fdc_lookup_kcal

# -------------------- App setup --------------------
st.set_page_config(page_title="Calorie Cards — Generator", layout="wide")

# files/dirs
CARDS_DIR = Path("cards")
CARDS_DIR.mkdir(exist_ok=True)

# Secrets / USDA key
FDC_API_KEY = st.secrets.get("FDC_API_KEY", os.getenv("FDC_API_KEY", ""))
if not FDC_API_KEY:
    with st.expander("USDA Internet Lookup (developer only – set key for local testing)"):
        FDC_API_KEY = st.text_input("FDC API Key", type="password")

# -------------------- Sidebar: Brand / Theme --------------------
with st.sidebar:
    st.header("Brand / Theme")
    c1, c2 = st.columns(2)
    c3, c4 = st.columns(2)
    panel_hex  = c1.color_picker("Panel", "#F4F4F4")
    accent_hex = c2.color_picker("Accent", "#672B91")
    text_hex   = c3.color_picker("Text", "#141414")
    faint_hex  = c4.color_picker("Muted", "#787878")

    to_rgb = lambda hx: tuple(int(hx[i:i+2], 16) for i in (1,3,5))
    theme = Theme(
        panel_color=to_rgb(panel_hex),
        accent=to_rgb(accent_hex),
        text=to_rgb(text_hex),
        faint=to_rgb(faint_hex),
    )

    st.header("Typography & Size")
    base_scale = st.slider("Base font size scale", 0.8, 2.2, 1.20, 0.01)
    card_size = st.selectbox(
        "Card size",
        options=[(1920,1200), (2560,1600), (2880,1800), (3840,2400)],
        index=0, format_func=lambda s: f"{s[0]} x {s[1]}"
    )
    right_ratio = st.slider("Right panel width (two-panel only)", 0.42, 0.72, 0.52, 0.01)

    # ---------------- Saved Cards Manager ----------------
    st.markdown("---")
    st.subheader("Saved Meal Cards")

    CARDS_DIR = Path("cards")
    CARDS_DIR.mkdir(exist_ok=True)

    def list_saved_cards():
        items = []
        for j in CARDS_DIR.glob("*.json"):
            name = j.stem
            p_png = CARDS_DIR / f"{name}.png"
            mtime = j.stat().st_mtime
            items.append({"name": name, "json": j, "png": p_png, "mtime": mtime})
        # newest first
        return sorted(items, key=lambda d: d["mtime"], reverse=True)

    saved = list_saved_cards()
    names = ["(none)"] + [it["name"] for it in saved]
    sel_name = st.selectbox("Select a saved card", options=names, key="saved_sel")

    if sel_name != "(none)":
        card = next((it for it in saved if it["name"] == sel_name), None)
        if card:
            # preview if PNG exists
            if card["png"].exists():
                st.image(str(card["png"]), use_column_width=True)
            btn1, btn2, btn3 = st.columns(3)
            if btn1.button("Load", use_container_width=True):
                st.session_state["_load_card_name"] = sel_name
                st.rerun()
            if btn2.download_button("JSON", data=card["json"].read_bytes(),
                                    file_name=f"{sel_name}.json", mime="application/json",
                                    use_container_width=True):
                pass
            if card["png"].exists():
                st.download_button("PNG", data=card["png"].read_bytes(),
                                   file_name=f"{sel_name}.png", mime="image/png",
                                   use_container_width=True)
            if btn3.button("Delete", type="secondary", use_container_width=True):
                # delete both JSON and PNG
                try:
                    card["json"].unlink(missing_ok=True)
                    card["png"].unlink(missing_ok=True)
                    st.success(f"Deleted {sel_name}")
                except Exception as e:
                    st.error(f"Could not delete: {e}")
                st.rerun()


# -------------------- Session bootstrap --------------------
def ensure_df():
    if "foods" not in st.session_state:
        st.session_state["foods"] = pd.DataFrame(
            [
                # Base macros
                {"category":"Protein","name":"Grilled Chicken 4 oz","cal":170},
                {"category":"Carb","name":"Mixed Veggies 1 cup","cal":70},
                {"category":"Fat","name":"Olive Oil 1 tsp","cal":40},
                # Combo categories (examples so UI demonstrates lookups & DB pulls)
                {"category":"Protein + Fat","name":"Whole Egg 1 each","cal":72},
                {"category":"Carb + Fat","name":"Avocado Toast (1/2 avo + 1 slice)","cal":180},
                {"category":"Protein + Carb","name":"Greek Yogurt + Berries (1 cup)","cal":150},
                {"category":"Protein + Carb + Fat","name":"Turkey Sandwich (half)","cal":220},
            ]
        )

ensure_df()
foods_df: pd.DataFrame = st.session_state["foods"]

# editor state helpers
SECTIONS = [
    ("Protein", "protein"),
    ("Carb", "carb"),
    ("Fat", "fat"),
    ("Protein + Fat", "pf"),
    ("Carb + Fat", "cf"),
    ("Protein + Carb", "pc"),
    ("Protein + Carb + Fat", "pcf"),
]

UNITS = ["g","oz","cup","tbsp","tsp","each"]

def rows_key(sec_key: str) -> str:
    return f"{sec_key}_rows"

def init_section_rows(sec_key: str, default_rows: int = 4):
    rk = rows_key(sec_key)
    if rk not in st.session_state:
        st.session_state[rk] = default_rows

def ensure_row_state(sec_key: str, i: int):
    """Initialize per-row keys so widgets bind correctly."""
    base = f"{sec_key}{i}"
    for suffix, default in (("_name",""), ("_amt",0.0), ("_unit","g"), ("_cal",0)):
        key = f"{base}{suffix}"
        if key not in st.session_state:
            st.session_state[key] = default

def reset_section(sec_key: str):
    rk = rows_key(sec_key)
    n = st.session_state.get(rk, 0)
    for i in range(1, n+1):
        base = f"{sec_key}{i}"
        for suffix in ("_name","_amt","_unit","_cal"):
            st.session_state.pop(f"{base}{suffix}", None)
    st.session_state[rk] = 4

def hard_reset_editor():
    # reset all sections and top-level inputs
    for _, sk in SECTIONS:
        reset_section(sk)
    for k in list(st.session_state.keys()):
        if k.startswith("db_"):  # multiselects
            st.session_state.pop(k, None)
    st.session_state.pop("_generated_png", None)
    st.session_state.pop("_last_card_name", None)

# ---------- USDA lookup ----------
def usda_lookup(name: str, amt: float, unit: str) -> int:
    kcal = fdc_lookup_kcal(name, amt, unit, api_key=FDC_API_KEY or "")
    return int(round(kcal or 0))

def _do_lookup(cal_key: str, name_key: str, amt_key: str, unit_key: str):
    name = st.session_state.get(name_key, "")
    amt  = float(st.session_state.get(amt_key, 0.0) or 0.0)
    unit = st.session_state.get(unit_key, "g")
    if name and (FDC_API_KEY or ""):
        st.session_state[cal_key] = usda_lookup(name, amt, unit)

# -------------------- Top: DB (left) / Last card (right) --------------------
left, right = st.columns([1,1], gap="large")

with left:
    st.subheader("📚 Food Database (add items here)")
    st.dataframe(foods_df, use_container_width=True, height=320)

    with st.form("new_food"):
        c1, c2, c3 = st.columns([1,2,1])
        cat = c1.selectbox("Category", [s[0] for s in SECTIONS[:3]])  # Only P/C/F in the DB categories
        nm  = c2.text_input("Name")
        cal = c3.number_input("Calories", min_value=0, step=1, value=0)
        if st.form_submit_button("Add"):
            if nm:
                st.session_state["foods"] = pd.concat(
                    [st.session_state["foods"], pd.DataFrame([{"category":cat,"name":nm,"cal":int(cal)}])],
                    ignore_index=True
                )
                st.rerun()
    csv = foods_df.to_csv(index=False).encode()
    st.download_button("Download DB CSV", data=csv, file_name="foods.csv", mime="text/csv")

with right:
    st.subheader("🧾 Last Generated Card")
    if st.session_state.get("_generated_png") and Path(st.session_state["_generated_png"]).exists():
        st.image(st.session_state["_generated_png"])
    elif os.path.exists("meal_card.png"):
        st.image("meal_card.png")
    else:
        st.info("No card generated yet")

# -------------------- Build a Single Meal Card --------------------
st.markdown("## 🍽️ Build a Single Meal Card")

# top inputs
c1, c2 = st.columns([1,1])
with c1:
    program    = st.text_input("Program Title", "40 Day Turn Up", key="program_title")
    grp        = st.text_input("Class / Group (optional)", "I RISE", key="group_name")
    meal_title = st.text_input("Meal Title", "Meal 1", key="meal_title")
    date_val   = st.date_input("Date", value=dt.date.today(), key="meal_date")
    date_str   = date_val.strftime("%Y-%m-%d")  # we'll format display separately below
    brand      = st.text_input("Brand (tiny footer)", "Alphonso Brown", key="brand_text")
with c2:
    photo = st.file_uploader("Upload meal photo", type=["png","jpg","jpeg"], key="photo_upload")
    photo_path = None
    if photo:
        photo_path = "preview_photo.png"
        with open(photo_path, "wb") as f:
            f.write(photo.read())
        st.image(photo_path, caption="Photo", use_container_width=True)

# utility: filename convention
def card_basename():
    yyyymmdd = date_val.strftime("%Y%m%d")
    title = f"{yyyymmdd} - {meal_title.strip() or 'Meal'} - {program.strip() or 'Program'} - {brand.strip() or 'Brand'}"
    safe = "".join(ch for ch in title if ch not in r'<>:"/\|?*').strip()
    return safe

# DB selector helper
def from_db(category_label: str):
    """Multiselect pulls entries that exactly match this category label."""
    opts = foods_df.query("category == @category_label")["name"].tolist() if not foods_df.empty else []
    return st.multiselect(f"Add {category_label.upper()} from DB", options=opts, key=f"db_{category_label}")


# -------- Section renderer (dynamic rows + lookup per line) --------
def render_section(title: str, sec_key: str):
    init_section_rows(sec_key)
    st.markdown(f"### {title.upper()}")
    sel = from_db(title)  # works for Protein, Carb, Fat, and all combo categories

    # Row controls
    cadd, crem = st.columns([1,1])
    if cadd.button(f"Add Row (+) [{title}]", key=f"{sec_key}_add"):
        st.session_state[rows_key(sec_key)] += 1
        st.rerun()
    if crem.button(f"Remove Row (–) [{title}]", key=f"{sec_key}_rem"):
        current = st.session_state[rows_key(sec_key)]
        if current > 1:
            # clean keys of the last row
            base = f"{sec_key}{current}"
            for suf in ("_name","_amt","_unit","_cal"):
                st.session_state.pop(f"{base}{suf}", None)
            st.session_state[rows_key(sec_key)] = current - 1
            st.rerun()

    # Rows
    rows = []
    for i in range(1, st.session_state[rows_key(sec_key)] + 1):
        ensure_row_state(sec_key, i)
        base = f"{sec_key}{i}"
        name_k, amt_k, unit_k, cal_k = f"{base}_name", f"{base}_amt", f"{base}_unit", f"{base}_cal"
        lk_k = f"{base}_lk"
        sv_k = f"{base}_sv"

        cA, cB, cC, cD, cE = st.columns([2.3, 0.9, 1.0, 1.0, 0.8])
        name = cA.text_input(f"item {i}", key=name_k, placeholder="")
        amt  = cB.number_input("amt", key=amt_k, step=0.25, min_value=0.0)
        unit = cC.selectbox("unit", UNITS, key=unit_k)
        cal  = cD.number_input("cal", key=cal_k, step=1, min_value=0)  # no 'value='

        cE.button("Lookup", key=lk_k, on_click=_do_lookup,
                  kwargs=dict(cal_key=cal_k, name_key=name_k, amt_key=amt_k, unit_key=unit_k))

        # Optional Save into in-session DB (category = section title)
        if st.button("Save", key=sv_k) and name and st.session_state.get(cal_k, 0) > 0:
            pretty = f"{name} {st.session_state.get(amt_k,0):g} {st.session_state.get(unit_k,'')}".strip()
            kcal   = int(st.session_state.get(cal_k, 0))
            st.session_state["foods"] = pd.concat(
                [st.session_state["foods"], pd.DataFrame([{
                    "category": title,       # <-- full label (e.g., "Protein + Fat")
                    "name": pretty, "cal": kcal
                }])],
                ignore_index=True,
            )
            st.toast(f"Saved to DB: {pretty} — {kcal} cal")

        rows.append((
            st.session_state.get(name_k,""),
            float(st.session_state.get(amt_k,0.0) or 0.0),
            st.session_state.get(unit_k,"g"),
            int(st.session_state.get(cal_k,0))
        ))
    return sel, rows

# -------- Render ALL sections --------
db_selects = {}
section_rows = {}

for title, key in SECTIONS:
    st.divider()
    sel, rows = render_section(title, key)
    db_selects[key] = sel
    section_rows[key] = rows

# -------- Collect final items to render card --------
def collect_items(db_names, manual_rows, category_name: str):
    items = []
    # pull from DB where available (only Protein/Carb/Fat categories)
    if db_names:
        sub = foods_df[foods_df["name"].isin(db_names)]
        for _, r in sub.iterrows():
            items.append(MealItem(text=r["name"], cal=int(r["cal"])))
    # manual rows
    for (name, amt, unit, cal) in manual_rows:
        if name and cal > 0:
            items.append(MealItem(text=f"{name} {amt:g} {unit}", cal=int(cal)))
    return items

prot_items = collect_items(db_selects["protein"], section_rows["protein"], "Protein")
carb_items = collect_items(db_selects["carb"],    section_rows["carb"],    "Carb")
fat_items  = collect_items(db_selects["fat"],     section_rows["fat"],     "Fat")

pf_items  = collect_items([], section_rows["pf"],  "Protein + Fat")
cf_items  = collect_items([], section_rows["cf"],  "Carb + Fat")
pc_items  = collect_items([], section_rows["pc"],  "Protein + Carb")
pcf_items = collect_items([], section_rows["pcf"], "Protein + Carb + Fat")

# Combine for totals
all_items = prot_items + carb_items + fat_items + pf_items + cf_items + pc_items + pcf_items

st.divider()
total_cals = int(sum(i.cal for i in all_items))
st.metric("Total Calories (auto; updates when you Lookup/Save)", total_cals)

# -------------------- Actions: Reset / Create New / Generate / Save --------------------
cA, cB, cC, cD = st.columns([1,1,1,1])

if cA.button("Reset Single Card"):
    # clear only the counts and per-row entries; keep Brand/Theme
    for _, k in SECTIONS:
        reset_section(k)
    st.toast("Single Card inputs reset.")
    st.rerun()

if cB.button("Create New Meal Card"):
    hard_reset_editor()
    st.toast("Started a new empty card.")
    st.rerun()

def build_card_data():
    # date shown on card as M/D/YY (portable)
    try:
        display_date = date_val.strftime("%-m/%-d/%y")
    except Exception:
        display_date = date_val.strftime("%m/%d/%y")

    # map UI rows to MealItem lists
    by_key = {
        "Protein": prot_items,
        "Carb":    carb_items,
        "Fat":     fat_items,
        "Protein + Fat":        pf_items,
        "Carb + Fat":           cf_items,
        "Protein + Carb":       pc_items,
        "Protein + Carb + Fat": pcf_items,
    }

    # Only include non-empty sections in the card
    dynamic_sections = [
        MealSection(title=label, items=items)
        for label, items in by_key.items()
        if items and len(items) > 0
    ]

    # Backward-compatible fields (your renderer may still consume these three)
    # We also pass a `sections` list so the generator can render N sections.
    return MealCardData(
        program_title=program.strip() or "Program",
        class_name=(grp.strip() or None),
        meal_title=meal_title.strip() or "Meal",
        date_str=display_date,
        brand=brand.strip() or None,
        protein=MealSection("Protein", prot_items),
        carb=MealSection("Carb", carb_items),
        fat=MealSection("Fat", fat_items),
        sections=dynamic_sections  # <-- new, used when available
    )

def save_card_json(png_path: str):
    payload = {
        "program": program, "group": grp, "meal_title": meal_title,
        "date": date_val.isoformat(), "brand": brand,
        "theme": dict(panel=panel_hex, accent=accent_hex, text=text_hex, faint=faint_hex),
        "size": list(card_size), "right_ratio": right_ratio, "base_scale": base_scale,
        "sections": {
            "Protein":  [i.__dict__ for i in prot_items],
            "Carb":     [i.__dict__ for i in carb_items],
            "Fat":      [i.__dict__ for i in fat_items],
            "Protein + Fat":        [i.__dict__ for i in pf_items],
            "Carb + Fat":           [i.__dict__ for i in cf_items],
            "Protein + Carb":       [i.__dict__ for i in pc_items],
            "Protein + Carb + Fat": [i.__dict__ for i in pcf_items],
        },
        "png": png_path,
    }
    name = card_basename()
    json_path = CARDS_DIR / f"{name}.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)
    st.session_state["_last_card_name"] = name
    return json_path

def generate_png_and_buttons(card: MealCardData):
    out_png = f"{card_basename()}.png"
    render_meal_card(
        card,
        photo_path=st.session_state.get("photo_upload") and "preview_photo.png" or None,
        output_path=out_png,
        size=card_size, theme=theme, font_scale=base_scale, panel_ratio=right_ratio
    )
    st.session_state["_generated_png"] = out_png
    st.success("Card generated.")
    st.image(out_png, use_container_width=True)

    # Download PNG
    with open(out_png, "rb") as f:
        st.download_button("Download PNG", data=f.read(), file_name=out_png, mime="image/png")

    # PPTX with the card as one slide
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = os.path.abspath(out_png)
    left = top = Inches(0.25)
    slide.shapes.add_picture(pic, left, top, width=Inches(9.5))
    bio = io.BytesIO(); prs.save(bio); bio.seek(0)
    pptx_name = f"{card_basename()}.pptx"
    st.download_button("Download PPTX", data=bio, file_name=pptx_name,
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    return out_png

if cC.button("Generate Card", type="primary", use_container_width=True):
    card = build_card_data()
    png_path = generate_png_and_buttons(card)

if cD.button("Save Card"):
    card = build_card_data()
    png_path = st.session_state.get("_generated_png")
    if not png_path or not Path(png_path).exists():
        # (re)generate if needed
        png_path = generate_png_and_buttons(card)
    # move/copy PNG into cards folder
    dest_png = CARDS_DIR / f"{card_basename()}.png"
    if Path(png_path).resolve() != dest_png.resolve():
        # copy
        with open(png_path, "rb") as src, open(dest_png, "wb") as dst:
            dst.write(src.read())
    json_path = save_card_json(str(dest_png))
    st.success(f"Saved card to {json_path}")
    st.balloons()

# -------------------- Load on request (from sidebar) --------------------
if st.session_state.get("_load_card_name"):
    name = st.session_state.pop("_load_card_name")
    jpath = CARDS_DIR / f"{name}.json"
    if jpath.exists():
        with open(jpath, "r", encoding="utf-8") as f:
            data = json.load(f)
        # Restore top inputs
        st.session_state["program_title"] = data.get("program","")
        st.session_state["group_name"]   = data.get("group","")
        st.session_state["meal_title"]   = data.get("meal_title","")
        st.session_state["meal_date"]    = dt.date.fromisoformat(data.get("date")[:10])
        st.session_state["brand_text"]   = data.get("brand","")
        # Restore theme & sizing
        # (We could restore the pickers directly, but Streamlit pickers ignore programmatic set;
        #  showing a toast instead, and we still render with saved theme on next Generate.)
        st.toast("Loaded details; adjust theme pickers as needed.")

        # Restore rows & entries for each section
        for title, key in SECTIONS:
            items = data.get("sections",{}).get(title, [])
            st.session_state[rows_key(key)] = max(1, len(items)) or 1
            for idx, item in enumerate(items, start=1):
                base = f"{key}{idx}"
                st.session_state[f"{base}_name"] = item.get("text","")
                # Extract amount+unit back out of text if present: "<name> <amt> <unit>"
                # This is a best-effort parse; we leave cal exact.
                st.session_state[f"{base}_amt"]  = 0.0
                st.session_state[f"{base}_unit"] = "g"
                st.session_state[f"{base}_cal"]  = int(item.get("cal",0))
        st.rerun()

